import os
from typing import Any, Dict, List

import pandas as pd
import pypdfium2 as pdfium
import pytesseract
from docx import Document
from PIL import Image, ImageOps, UnidentifiedImageError
from pptx import Presentation
from pypdf import PdfReader
from pytesseract import TesseractNotFoundError


SUPPORTED_EXTENSIONS = [
    ".pdf",
    ".txt",
    ".docx",
    ".pptx",
    ".csv",
    ".xlsx",
    ".jpg",
    ".jpeg",
    ".png",
]

SourceSection = Dict[str, Any]

OCR_LANGUAGE = os.getenv("OCR_LANGUAGE", "eng")
OCR_DPI = int(os.getenv("OCR_DPI", "180"))
OCR_TIMEOUT_SECONDS = int(os.getenv("OCR_TIMEOUT_SECONDS", "15"))
OCR_NATIVE_TEXT_MIN_CHARS = int(
    os.getenv("OCR_NATIVE_TEXT_MIN_CHARS", "40")
)
MAX_OCR_PAGES = int(os.getenv("MAX_OCR_PAGES", "10"))
MAX_IMAGE_PIXELS = int(os.getenv("MAX_IMAGE_PIXELS", "25000000"))


class DocumentExtractionError(Exception):
    """Raised when a supported file cannot be read or processed."""


class OCRUnavailableError(DocumentExtractionError):
    """Raised when OCR is required but Tesseract is unavailable."""


def get_file_extension(file_name: str) -> str:
    return os.path.splitext(file_name)[1].lower()


def is_supported_file(file_name: str) -> bool:
    return get_file_extension(file_name) in SUPPORTED_EXTENSIONS


def _clean_text(value: Any) -> str:
    if value is None:
        return ""

    text = str(value).replace("\x00", "")
    lines = [line.rstrip() for line in text.splitlines()]
    return "\n".join(lines).strip()


def _readable_character_count(text: str) -> int:
    return sum(character.isalnum() for character in text)


def _has_enough_native_text(text: str) -> bool:
    return _readable_character_count(text) >= OCR_NATIVE_TEXT_MIN_CHARS


def _validate_image_size(image: Image.Image) -> None:
    width, height = image.size
    pixels = width * height

    if pixels <= 0:
        raise DocumentExtractionError("The image has invalid dimensions.")

    if pixels > MAX_IMAGE_PIXELS:
        raise DocumentExtractionError(
            "The image is too large to process safely."
        )


def _prepare_image_for_ocr(image: Image.Image) -> Image.Image:
    image = ImageOps.exif_transpose(image)
    _validate_image_size(image)

    if image.mode not in ("RGB", "L"):
        image = image.convert("RGB")

    return image


def _run_ocr(image: Image.Image) -> str:
    image = _prepare_image_for_ocr(image)

    try:
        text = pytesseract.image_to_string(
            image,
            lang=OCR_LANGUAGE,
            timeout=OCR_TIMEOUT_SECONDS,
            config="--psm 6",
        )
    except TesseractNotFoundError as error:
        raise OCRUnavailableError(
            "OCR is not available on the server because Tesseract "
            "is not installed or cannot be found."
        ) from error
    except RuntimeError as error:
        raise DocumentExtractionError(
            "OCR processing timed out for this page or image."
        ) from error
    except Exception as error:
        raise DocumentExtractionError(
            "OCR could not process this page or image."
        ) from error

    return _clean_text(text)


def _render_pdf_page_for_ocr(
    pdf_document: pdfium.PdfDocument,
    page_index: int,
) -> Image.Image:
    try:
        page = pdf_document[page_index]
        scale = max(OCR_DPI / 72.0, 1.0)
        bitmap = page.render(scale=scale)
        image = bitmap.to_pil()
        return _prepare_image_for_ocr(image)
    except DocumentExtractionError:
        raise
    except Exception as error:
        raise DocumentExtractionError(
            f"Could not render PDF page {page_index + 1} for OCR."
        ) from error


def _cell_to_text(value: Any) -> str:
    if value is None:
        return ""

    try:
        if pd.isna(value):
            return ""
    except (TypeError, ValueError):
        pass

    return str(value).strip()


def _dataframe_rows(dataframe: pd.DataFrame) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    columns = [str(column).strip() for column in dataframe.columns]

    for row_number, (_, record) in enumerate(
        dataframe.iterrows(),
        start=2,
    ):
        values = []

        for column in columns:
            value = _cell_to_text(record.get(column, ""))
            if value:
                values.append(f"{column}: {value}")

        if values:
            rows.append(
                {
                    "text": " | ".join(values),
                    "row_number": row_number,
                }
            )

    return rows


def extract_text_from_pdf(file_path: str) -> List[SourceSection]:
    """
    Hybrid PDF extraction:
    - use embedded/selectable text when available;
    - run OCR only on pages with little or no native text.
    """
    try:
        reader = PdfReader(file_path)
    except Exception as error:
        raise DocumentExtractionError(
            "The PDF could not be opened or is corrupted."
        ) from error

    sections: List[SourceSection] = []
    pdf_for_ocr = None
    ocr_pages_used = 0

    try:
        for page_number, page in enumerate(reader.pages, start=1):
            try:
                native_text = _clean_text(page.extract_text() or "")
            except Exception:
                native_text = ""

            if _has_enough_native_text(native_text):
                sections.append(
                    {
                        "text": native_text,
                        "metadata": {
                            "source_type": "pdf",
                            "page_number": page_number,
                            "source_label": f"Page {page_number}",
                            "extraction_method": "native",
                            "ocr_used": False,
                        },
                    }
                )
                continue

            if ocr_pages_used >= MAX_OCR_PAGES:
                if native_text:
                    sections.append(
                        {
                            "text": native_text,
                            "metadata": {
                                "source_type": "pdf",
                                "page_number": page_number,
                                "source_label": f"Page {page_number}",
                                "extraction_method": "native",
                                "ocr_used": False,
                            },
                        }
                    )
                continue

            if pdf_for_ocr is None:
                try:
                    pdf_for_ocr = pdfium.PdfDocument(file_path)
                except Exception as error:
                    raise DocumentExtractionError(
                        "The PDF could not be rendered for OCR."
                    ) from error

            image = _render_pdf_page_for_ocr(
                pdf_document=pdf_for_ocr,
                page_index=page_number - 1,
            )
            ocr_text = _run_ocr(image)
            ocr_pages_used += 1

            final_text = ocr_text or native_text

            if not final_text:
                continue

            sections.append(
                {
                    "text": final_text,
                    "metadata": {
                        "source_type": "pdf",
                        "page_number": page_number,
                        "source_label": f"Page {page_number}",
                        "extraction_method": (
                            "ocr" if ocr_text else "native"
                        ),
                        "ocr_used": bool(ocr_text),
                    },
                }
            )
    finally:
        if pdf_for_ocr is not None:
            try:
                pdf_for_ocr.close()
            except Exception:
                pass

    return sections


def extract_text_from_image(file_path: str) -> List[SourceSection]:
    try:
        with Image.open(file_path) as image:
            image.load()
            prepared_image = _prepare_image_for_ocr(image.copy())
    except (UnidentifiedImageError, OSError) as error:
        raise DocumentExtractionError(
            "The image could not be opened or is corrupted."
        ) from error
    except Image.DecompressionBombError as error:
        raise DocumentExtractionError(
            "The image is too large to process safely."
        ) from error

    text = _run_ocr(prepared_image)

    if not text:
        return []

    return [
        {
            "text": text,
            "metadata": {
                "source_type": "image",
                "image_number": 1,
                "source_label": "Image",
                "extraction_method": "ocr",
                "ocr_used": True,
            },
        }
    ]


def extract_text_from_txt(file_path: str) -> List[SourceSection]:
    try:
        with open(
            file_path,
            "r",
            encoding="utf-8",
            errors="ignore",
        ) as file:
            raw_text = file.read()
    except OSError as error:
        raise DocumentExtractionError(
            "The text file could not be read."
        ) from error

    lines: List[SourceSection] = []

    for line_number, line in enumerate(raw_text.splitlines(), start=1):
        line_text = _clean_text(line)

        if line_text:
            lines.append(
                {
                    "text": line_text,
                    "metadata": {
                        "source_type": "txt",
                        "line_number": line_number,
                    },
                }
            )

    if not lines:
        text = _clean_text(raw_text)

        if text:
            lines.append(
                {
                    "text": text,
                    "metadata": {
                        "source_type": "txt",
                        "line_number": 1,
                    },
                }
            )

    return lines


def extract_text_from_docx(file_path: str) -> List[SourceSection]:
    try:
        document = Document(file_path)
    except Exception as error:
        raise DocumentExtractionError(
            "The Word document could not be opened or is corrupted."
        ) from error

    sections: List[SourceSection] = []
    paragraph_number = 0

    for paragraph in document.paragraphs:
        paragraph_text = _clean_text(paragraph.text)

        if not paragraph_text:
            continue

        paragraph_number += 1
        sections.append(
            {
                "text": paragraph_text,
                "metadata": {
                    "source_type": "docx",
                    "paragraph_number": paragraph_number,
                },
            }
        )

    for table_number, table in enumerate(document.tables, start=1):
        for row_number, row in enumerate(table.rows, start=1):
            cells = []

            for cell in row.cells:
                cell_text = _clean_text(cell.text)
                if cell_text:
                    cells.append(cell_text)

            if cells:
                sections.append(
                    {
                        "text": " | ".join(cells),
                        "metadata": {
                            "source_type": "docx",
                            "table_number": table_number,
                            "table_row_number": row_number,
                        },
                    }
                )

    return sections


def extract_text_from_pptx(file_path: str) -> List[SourceSection]:
    try:
        presentation = Presentation(file_path)
    except Exception as error:
        raise DocumentExtractionError(
            "The PowerPoint file could not be opened or is corrupted."
        ) from error

    slides: List[SourceSection] = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_parts: List[str] = []

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape_text = _clean_text(shape.text)
                if shape_text:
                    slide_parts.append(shape_text)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = []

                    for cell in row.cells:
                        cell_text = _clean_text(cell.text)
                        if cell_text:
                            cells.append(cell_text)

                    if cells:
                        slide_parts.append(" | ".join(cells))

        slide_text = "\n".join(slide_parts).strip()

        if slide_text:
            slides.append(
                {
                    "text": slide_text,
                    "metadata": {
                        "source_type": "pptx",
                        "slide_number": slide_number,
                        "source_label": f"Slide {slide_number}",
                        "extraction_method": "native",
                        "ocr_used": False,
                    },
                }
            )

    return slides


def extract_text_from_csv(file_path: str) -> List[SourceSection]:
    try:
        dataframe = pd.read_csv(
            file_path,
            dtype=str,
            keep_default_na=False,
        )
    except Exception as error:
        raise DocumentExtractionError(
            "The CSV file could not be read."
        ) from error

    if dataframe.empty and len(dataframe.columns) == 0:
        return []

    rows = _dataframe_rows(dataframe)
    header = "Columns: " + " | ".join(
        str(column) for column in dataframe.columns
    )
    text_parts = [header]
    text_parts.extend(row["text"] for row in rows)

    return [
        {
            "text": "\n".join(text_parts).strip(),
            "rows": rows,
            "metadata": {
                "source_type": "csv",
                "sheet_name": "CSV",
                "source_label": "CSV data",
                "extraction_method": "native",
                "ocr_used": False,
            },
        }
    ]


def extract_text_from_excel(file_path: str) -> List[SourceSection]:
    try:
        excel_file = pd.ExcelFile(file_path)
    except Exception as error:
        raise DocumentExtractionError(
            "The Excel file could not be opened or is corrupted."
        ) from error

    sheets: List[SourceSection] = []

    for sheet_name in excel_file.sheet_names:
        try:
            dataframe = pd.read_excel(
                file_path,
                sheet_name=sheet_name,
                dtype=str,
                keep_default_na=False,
            )
        except Exception as error:
            raise DocumentExtractionError(
                f'Excel sheet "{sheet_name}" could not be read.'
            ) from error

        if dataframe.empty and len(dataframe.columns) == 0:
            continue

        rows = _dataframe_rows(dataframe)
        header = "Columns: " + " | ".join(
            str(column) for column in dataframe.columns
        )
        text_parts = [header]
        text_parts.extend(row["text"] for row in rows)

        sheets.append(
            {
                "text": "\n".join(text_parts).strip(),
                "rows": rows,
                "metadata": {
                    "source_type": "xlsx",
                    "sheet_name": str(sheet_name),
                    "source_label": f'Sheet "{sheet_name}"',
                    "extraction_method": "native",
                    "ocr_used": False,
                },
            }
        )

    return sheets


def extract_text_from_document(
    file_path: str,
    file_name: str,
) -> List[SourceSection]:
    file_extension = get_file_extension(file_name)

    if file_extension == ".pdf":
        return extract_text_from_pdf(file_path)

    if file_extension == ".txt":
        return extract_text_from_txt(file_path)

    if file_extension == ".docx":
        return extract_text_from_docx(file_path)

    if file_extension == ".pptx":
        return extract_text_from_pptx(file_path)

    if file_extension == ".csv":
        return extract_text_from_csv(file_path)

    if file_extension == ".xlsx":
        return extract_text_from_excel(file_path)

    if file_extension in {".jpg", ".jpeg", ".png"}:
        return extract_text_from_image(file_path)

    raise DocumentExtractionError(
        f"Unsupported file type: {file_extension}"
    )